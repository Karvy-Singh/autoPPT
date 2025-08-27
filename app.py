from __future__ import annotations

import io
import json
import hashlib
from dataclasses import dataclass
from typing import List, Dict, Optional

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor

st.set_page_config(page_title="Your Text, Your Style — PPTX Generator", layout="wide")

HIDE_STREAMLIT_STYLE = """
    <style>
        /* Cleaner, app-like chrome */
        header {visibility: hidden;}
        #MainMenu {visibility:hidden;}
        footer {visibility: hidden;}
        .stApp {background: #0b0f19;}
        .app-card {background:#101726; border:1px solid #1e2a44; padding:1.25rem; border-radius:16px;}
        .accent {color:#bcd7ff;}
        .muted {color:#95a3bd;}
        .cta {background:linear-gradient(135deg,#3a7bd5,#00d2ff); -webkit-background-clip:text; -webkit-text-fill-color:transparent;}
        .small {font-size:0.9rem;}
        .tag {display:inline-block; border:1px solid #2a3656; padding:.15rem .5rem; border-radius:6px; margin-right:.35rem; color:#b5c7f2; font-size:.8rem}
        .success {color:#7ee787;}
        .warn {color:#ffd480;}
    </style>
"""
st.markdown(HIDE_STREAMLIT_STYLE, unsafe_allow_html=True)

st.markdown(
    """
    <div class="app-card">
        <h1 style="margin-top:0">Your Text, <span class="cta">Your Style</span> — Auto‑Generate a Presentation</h1>
        <p class="muted">Paste text or markdown, pick a template (.pptx/.potx), and get a matching deck. The app
        uses your template’s layouts, theme, and images — it does <em>not</em> generate new images.</p>
        <div>
            <span class="tag">LLM‑assisted structure</span>
            <span class="tag">Template‑true styling</span>
            <span class="tag">Downloads .pptx</span>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

@dataclass
class SlideSpec:
    title: str
    bullets: List[str]
    layout: Optional[str] = None  # "title_and_content", "title_only", "section_header", "two_content"
    notes: Optional[str] = None
    image_hint: Optional[str] = None

@dataclass
class Outline:
    slides: List[SlideSpec]

SYSTEM_PROMPT = (
    "You convert long text/markdown into a clean presentation outline. "
    "Return STRICT, VALID JSON only with this shape: {\"slides\":[{"
    "\"title\": str, \"bullets\": [str,...], \"layout\": (\"title_and_content\"|\"title_only\"|\"section_header\"|\"two_content\")?, "
    "\"notes\"?: str, \"image_hint\"?: str}]}. "
    "Choose a reasonable number of slides. Keep bullets concise (<=90 chars)."
)

USER_PROMPT_TEMPLATE = (
    "TASK: Convert the INPUT into slides. If markdown headings exist, use them as section breaks.\n"
    "GUIDANCE (may be empty): {guidance}\n"
    "CONSTRAINTS: Aim for ~{target_slides} slides when appropriate. Always return valid JSON only.\n"
    "INPUT BEGIN\n{body}\nINPUT END"
)


def call_openai(api_key: str, text: str, guidance: str, target_slides: int) -> str:
    try:
        import openai
        client = openai.OpenAI(api_key=api_key)
        rsp = client.chat.completions.create(
            model="gpt-4o-mini-2024-07-18",
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {
                    "role": "user",
                    "content": USER_PROMPT_TEMPLATE.format(
                        guidance=guidance or "", body=text, target_slides=target_slides,
                    ),
                },
            ],
            temperature=0.2,
            response_format={"type": "json_object"},
        )
        return rsp.choices[0].message.content

    except Exception as e:
        raise RuntimeError(f"OpenAI error: {e}")


def call_anthropic(api_key: str, text: str, guidance: str, target_slides: int) -> str:
    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key)
        rsp = client.messages.create(
            model="claude-3-5-sonnet-20240620",
            max_tokens=4000,
            temperature=0.2,
            system=SYSTEM_PROMPT,
            messages=[
                {
                    "role": "user",
                    "content": USER_PROMPT_TEMPLATE.format(
                        guidance=guidance or "", body=text, target_slides=target_slides
                    ),
                }
            ],
        )
        # Anthropic content parts -> text
        parts = []
        for block in rsp.content:
            if block.type == "text":
                parts.append(block.text)
        return "\n".join(parts)
    except Exception as e:
        raise RuntimeError(f"Anthropic error: {e}")


def call_gemini(api_key: str, text: str, guidance: str, target_slides: int) -> str:
    try:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("gemini-1.5-pro-latest")
        prompt = (
            SYSTEM_PROMPT
            + "\n\n"
            + USER_PROMPT_TEMPLATE.format(
                guidance=guidance or "", body=text, target_slides=target_slides
            )
        )
        rsp = model.generate_content(prompt)
        return rsp.text
    except Exception as e:
        raise RuntimeError(f"Gemini error: {e}")


def outline_via_llm(provider: str, api_key: str, text: str, guidance: str, target_slides: int) -> Outline:
    if not text or len(text.strip()) == 0:
        raise ValueError("Please provide some input text.")

    if provider == "OpenAI":
        raw = call_openai(api_key, text, guidance, target_slides)
    elif provider == "Anthropic":
        raw = call_anthropic(api_key, text, guidance, target_slides)
    elif provider == "Gemini":
        raw = call_gemini(api_key, text, guidance, target_slides)
    elif provider == "Local heuristic":
        return heuristic_outline(text, guidance, target_slides)
    else:
        raise ValueError("Unsupported provider.")

    data = coerce_json(raw)
    slides = [SlideSpec(**s) for s in data.get("slides", [])]
    return Outline(slides=slides)

def heuristic_outline(text: str, guidance: str, target_slides: int) -> Outline:
    """A simple markdown-aware splitter in case no API key is used."""
    lines = [ln.strip() for ln in text.splitlines()]
    blocks: List[List[str]] = []
    cur: List[str] = []
    for ln in lines:
        if ln.startswith("#") and cur:
            blocks.append(cur)
            cur = [ln]
        else:
            cur.append(ln)
    if cur:
        blocks.append(cur)

    slides: List[SlideSpec] = []
    for block in blocks:
        if not block:
            continue
        title = block[0].lstrip("# ") if block[0].startswith("#") else (block[0][:80] + ("…" if len(block[0]) > 80 else ""))
        content = block[1:] if block[0].startswith("#") else block
        bullets = squash_to_bullets("\n".join(content))[:8]
        layout = "title_and_content" if bullets else "title_only"
        slides.append(SlideSpec(title=title or "Slide", bullets=bullets, layout=layout))

    # Thin or dense inputs
    if not slides:
        bullets = squash_to_bullets(text)[:8]
        slides = [SlideSpec(title="Overview", bullets=bullets, layout="title_and_content")]

    if target_slides and len(slides) > target_slides:
        slides = slides[:target_slides]

    return Outline(slides=slides)


def squash_to_bullets(text: str) -> List[str]:
    # Split on markdown bullets / lines / punctuation-ish.
    candidates = []
    for ln in text.splitlines():
        ln = ln.strip(" •-\t")
        if not ln:
            continue
        # further chop long lines
        parts = [p.strip() for p in reflow(ln, hard_max=90)]
        candidates.extend(parts)
    # de-dup nearby similar lines
    comp = []
    seen = set()
    for c in candidates:
        key = c.lower()
        if key not in seen:
            comp.append(c)
            seen.add(key)
    return comp


def reflow(s: str, hard_max: int = 90) -> List[str]:
    if len(s) <= hard_max:
        return [s]
    words = s.split()
    lines = []
    cur = []
    cur_len = 0
    for w in words:
        if cur_len + len(w) + (1 if cur else 0) > hard_max:
            lines.append(" ".join(cur))
            cur = [w]
            cur_len = len(w)
        else:
            cur.append(w)
            cur_len += len(w) + (1 if cur_len else 0)
    if cur:
        lines.append(" ".join(cur))
    return lines

def coerce_json(raw: str) -> Dict:
    """Try hard to extract valid JSON from an LLM response."""
    try:
        return json.loads(raw)
    except Exception:
        pass
    # find first { .. last }
    try:
        first = raw.find("{")
        last = raw.rfind("}")
        if first != -1 and last != -1:
            return json.loads(raw[first : last + 1])
    except Exception:
        pass
    # Minimal fallback
    return {"slides": []}


from typing import Tuple
import re


def extract_template_images(prs: Presentation) -> List[bytes]:
    imgs: List[bytes] = []
    hashes = set()

    def add_blob(b: bytes):
        h = hashlib.sha256(b).hexdigest()
        if h not in hashes:
            hashes.add(h)
            imgs.append(b)

    # From existing slides
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                add_blob(sh.image.blob)

    # From masters and layouts
    for master in prs.slide_masters:
        for sh in master.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                add_blob(sh.image.blob)
        for layout in master.slide_layouts:
            for sh in layout.shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    add_blob(sh.image.blob)

    return imgs


def clear_all_slides(prs: Presentation) -> None:
    """Remove all slides from a Presentation while keeping theme/layouts."""
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    if sldIdLst is None:
        return
    ids = list(sldIdLst)
    for sld in ids:
        sldIdLst.remove(sld)


def find_best_layout(prs: Presentation, hint: Optional[str]) -> Tuple[object, str]:
    """Pick a slide layout by name heuristic. Returns (layout, matched_name)."""
    target = (hint or "title_and_content").lower()

    def by_name(name_parts: List[str]):
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                name = (layout.name or "").lower()
                if all(p in name for p in name_parts):
                    return layout, layout.name
        return None

    # Mappings
    mapping = {
        "section_header": ["section"],
        "title_only": ["title", "only"],
        "two_content": ["two", "content"],
        "title_and_content": ["title", "content"],
    }
    if target in mapping:
        found = by_name(mapping[target])
        if found:
            return found

    # Common alternates
    for parts in (["title slide"], ["comparison"], ["blank"], ["content"]):
        found = by_name(parts)
        if found:
            return found

    # Fallback to first available layout
    for master in prs.slide_masters:
        if master.slide_layouts:
            layout = master.slide_layouts[0]
            return layout, layout.name

    # Should not happen
    raise RuntimeError("No slide layouts found in template.")


def fill_title(shape_collection, text: str):
    try:
        title_shape = shape_collection.title
        if title_shape:
            title_shape.text_frame.clear()
            title_shape.text = text
            return True
    except Exception:
        pass
    # fall back: first title placeholder
    for sh in shape_collection.placeholders:
        try:
            if sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                sh.text = text
                return True
        except Exception:
            continue
    return False


def fill_body(shape_collection, bullets: List[str]):
    # Find a body/content placeholder
    target = None
    for sh in shape_collection.placeholders:
        try:
            if sh.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT):
                target = sh
                break
        except Exception:
            continue
    if not target:
        return False

    tf = target.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
    return True


def maybe_fill_picture_placeholders(slide, image_pool: List[bytes]):
    if not image_pool:
        return
    # Round-robin image selection via session counter
    idx = st.session_state.get("_img_rr", 0)

    # find picture placeholders
    for sh in slide.shapes:
        try:
            if getattr(sh, "is_placeholder", False) and \
               sh.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                b = image_pool[idx % len(image_pool)]
                idx += 1
                from io import BytesIO
                sh.insert_picture(BytesIO(b))
        except Exception:
            continue

    st.session_state["_img_rr"] = idx


def add_freeform_picture(slide, image_pool: List[bytes]):
    if not image_pool:
        return
    # place bottom-right as a tasteful accent if no placeholder exists
    try:
        b = image_pool[st.session_state.get("_img_rr2", 0) % len(image_pool)]
        st.session_state["_img_rr2"] = st.session_state.get("_img_rr2", 0) + 1
        from io import BytesIO
        pic = slide.shapes.add_picture(BytesIO(b), Inches(7.0), Inches(4.5), width=Inches(2.5))
        # No explicit styling; preserve theme background
    except Exception:
        pass


def build_presentation(template_bytes: bytes, outline: Outline) -> bytes:
    # Load template as base so we inherit theme/fonts/layouts
    prs = Presentation(io.BytesIO(template_bytes))

    # Collect images BEFORE clearing (so we keep ones from sample slides too)
    image_pool = extract_template_images(prs)

    # Clear existing slides so only generated slides remain
    clear_all_slides(prs)

    # Build slides
    for ix, spec in enumerate(outline.slides):
        layout, matched = find_best_layout(prs, spec.layout)
        slide = prs.slides.add_slide(layout)

        # Title
        if spec.title:
            fill_title(slide.shapes, spec.title)

        # Content
        if spec.bullets:
            body_ok = fill_body(slide.shapes, spec.bullets)
        else:
            body_ok = False

        # Images: prefer picture placeholders, else freeform accent (optional)
        maybe_fill_picture_placeholders(slide, image_pool)
        if not body_ok:
            # If an otherwise-empty slide exists and we still have images, add one for balance
            add_freeform_picture(slide, image_pool)

        # Speaker notes
        if spec.notes:
            try:
                notes = slide.notes_slide
                notes.notes_text_frame.clear()
                notes.notes_text_frame.text = spec.notes
            except Exception:
                pass

    # Serialize
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

with st.sidebar:
    st.markdown("### Settings")

    provider = st.selectbox(
        "LLM Provider",
        ["OpenAI", "Anthropic", "Gemini", "Local heuristic"],
        help="Use your own API key. Choose 'Local heuristic' to avoid any external calls.",
    )

    api_key = st.text_input(
        "API key (never stored or logged)",
        type="password",
        placeholder="sk-… / anth-… / AIza…",
        help="Key is kept in memory only for this session and request.",
        disabled=(provider == "Local heuristic"),
    )

    target_slides = st.slider(
        "Target slide count (hint only)", 4, 40, 12,
        help="A hint for the LLM/heuristic — final count may vary based on your content.",
    )

    st.markdown("—" * 20)
    st.caption("Tips: Use a .potx or a .pptx with the exact theme & layouts you want.")

col_left, col_right = st.columns([7, 5], gap="large")

with col_left:
    st.markdown("#### 1) Paste your text or markdown")
    input_text = st.text_area(
        label="",
        height=320,
        placeholder="Paste anything — long-form prose, notes, or markdown…",
    )

    st.markdown("#### 2) Optional one‑line guidance")
    guidance = st.text_input(
        label="",
        placeholder="e.g. turn into an investor pitch deck / team update / lesson plan",
    )

    st.markdown("#### 3) Upload your PowerPoint template or presentation (.pptx/.potx)")
    upl = st.file_uploader(
        label="",
        type=["pptx", "potx"],
        accept_multiple_files=False,
        help="Your final deck will inherit this file's theme, layouts, fonts, and images.",
    )

with col_right:
    st.markdown("#### Preview & Generate")
    st.write(
        "Paste your content, choose a provider, upload a template, then click **Generate**. \n"
        "The app will outline slides, reuse template images where it can, and export a themed .pptx."
    )

    disabled = not (input_text and upl and (provider == "Local heuristic" or api_key))

    if st.button("Generate presentation", type="primary", disabled=disabled, use_container_width=True):
        if not upl:
            st.error("Please upload a .pptx or .potx template.")
        else:
            with st.spinner("Structuring content and composing slides…"):
                try:
                    outline = outline_via_llm(provider, api_key or "", input_text, guidance, target_slides)
                    if not outline.slides:
                        st.warning("No slides were produced — falling back to a simple overview.")
                        outline = heuristic_outline(input_text, guidance, target_slides)

                    pptx_bytes = build_presentation(upl.read(), outline)

                    st.success("Done! Your themed presentation is ready.")
                    st.download_button(
                        "Download .pptx",
                        data=pptx_bytes,
                        file_name="your-text-your-style.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )

                    # Simple preview facts
                    st.markdown(
                        f"<div class='app-card'><span class='success'>Generated</span> <b>{len(outline.slides)}</b> slides. "
                        f"Used template images: <b>{'Yes' if extract_template_images(Presentation(io.BytesIO(upl.getvalue()))) else 'No'}</b>.</div>",
                        unsafe_allow_html=True,
                    )
                except Exception as e:
                    st.error(f"Something went wrong: {e}")


if False:  # Set True while developing locally
    TEST_MD = """
# Vision & Mission
We build privacy-first AI.

# Problem
Teams waste hours formatting slides. Content is there, but design is hard.

# Solution
A tool that turns long text into a themed deck. Works with your PPT template.

# Traction
• 1k+ beta users  
• 200+ companies  
• 99.95% uptime

# Roadmap
Q3: enterprise SSO  
Q4: audit logs

# Ask
Pilot with 5 design partners.
"""
    

